"""그룹 자료 본문 추출.

Spring 이 허용하는 확장자 중 텍스트를 뽑을 수 있는 것만 다룬다.
이미지와 zip 은 색인 대상이 아니다.
"""

from __future__ import annotations

import csv
import io
import logging

log = logging.getLogger(__name__)

TEXT_EXTENSIONS = {"txt", "csv", "pdf", "docx", "xlsx", "pptx"}


class UnsupportedFormat(RuntimeError):
    pass


def extension_of(filename: str | None) -> str:
    if not filename or "." not in filename:
        return ""
    return filename.rsplit(".", 1)[1].lower()


def can_extract(filename: str | None) -> bool:
    return extension_of(filename) in TEXT_EXTENSIONS


def extract(content: bytes, filename: str | None) -> str:
    extension = extension_of(filename)
    if extension == "txt":
        return _decode(content)
    if extension == "csv":
        return _from_csv(content)
    if extension == "pdf":
        return _from_pdf(content)
    if extension == "docx":
        return _from_docx(content)
    if extension == "xlsx":
        return _from_xlsx(content)
    if extension == "pptx":
        return _from_pptx(content)
    raise UnsupportedFormat(f"본문을 추출할 수 없는 형식입니다: {extension or '확장자 없음'}")


def _decode(content: bytes) -> str:
    for encoding in ("utf-8", "cp949", "euc-kr"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _from_csv(content: bytes) -> str:
    """행을 '열이름: 값' 형태로 편다. 헤더만 반복되는 청크가 검색을 망치지 않게 한다."""
    text = _decode(content)
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return ""
    header, *body = rows
    if not body:
        return ", ".join(header)
    lines = []
    for row in body:
        pairs = [f"{name}: {value}" for name, value in zip(header, row) if value != ""]
        lines.append(", ".join(pairs))
    return "\n".join(lines)


def _from_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            pages.append(page.extract_text() or "")
        except Exception:  # 손상된 페이지 하나가 문서 전체를 버리게 두지 않는다.
            log.warning("PDF %d 페이지 추출 실패", index)
    return "\n\n".join(part for part in pages if part.strip())


def _from_docx(content: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(content))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _from_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    for sheet in workbook.worksheets:
        parts.append(f"[시트: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                parts.append(" | ".join(values))
    workbook.close()
    return "\n".join(parts)


def _from_pptx(content: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[슬라이드 {index}]")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)
